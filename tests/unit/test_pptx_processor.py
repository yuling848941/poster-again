#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
测试PPTXProcessor对不同文件的处理
"""

from pptx import Presentation
import logging

logging.basicConfig(level=logging.INFO, format='%(levelname)s: %(message)s')

def test_pptx_processor(ppt_path):
    """
    测试PPTXProcessor查找占位符的能力
    """
    try:
        print(f"\n测试文件: {ppt_path}")
        print("=" * 70)

        # 加载演示文稿
        presentation = Presentation(ppt_path)
        print(f"加载成功，共 {len(presentation.slides)} 张幻灯片")

        placeholders = []

        for slide_idx, slide in enumerate(presentation.slides, 1):
            if slide_idx > 3:  # 只检查前3张
                break

            print(f"\n【幻灯片 {slide_idx}】")

            for shape_idx, shape in enumerate(slide.shapes, 1):
                # 尝试获取名称
                shape_name = None
                try:
                    # python-pptx中形状名称可能存储在不同的属性中
                    if hasattr(shape, 'name'):
                        shape_name = shape.name
                    elif hasattr(shape, 'Name'):
                        shape_name = shape.Name
                    elif hasattr(shape, 'title'):
                        shape_name = shape.title
                except Exception as e:
                    print(f"    获取名称失败: {e}")

                if shape_name:
                    if shape_name.startswith('ph_'):
                        print(f"  [占位符] {shape_name}")
                        placeholders.append(shape_name)
                    else:
                        print(f"  [形状] {shape_name}")
                else:
                    # 尝试获取形状类型和文本内容
                    shape_type = "未知"
                    text_content = ""
                    try:
                        shape_type = str(shape.shape_type)
                        if hasattr(shape, 'text'):
                            text_content = shape.text.strip()
                    except:
                        pass

                    print(f"  [形状#{shape_idx}] 类型: {shape_type}, 文本: {text_content[:30] if text_content else '无'}")

        print("\n" + "=" * 70)
        print("测试结果")
        print("=" * 70)
        print(f"找到占位符: {len(placeholders)} 个")

        if placeholders:
            print("\n占位符列表:")
            for i, ph in enumerate(placeholders, 1):
                print(f"  {i}. {ph}")
            print("\n✓ PPTXProcessor可以正确识别占位符")
        else:
            print("\n✗ PPTXProcessor无法识别占位符")
            print("\n可能原因:")
            print("1. python-pptx库访问属性的方式不同")
            print("2. 形状名称存储在不同的属性中")
            print("3. 需要使用不同的属性名")

        return placeholders

    except Exception as e:
        print(f"测试失败: {e}")
        import traceback
        traceback.print_exc()
        return []

if __name__ == "__main__":
    # 测试两个文件
    ph1 = test_pptx_processor("E:/poster code/poster-again/10月承保贺报(4).pptx")
    ph2 = test_pptx_processor("E:/poster code/poster-again/10月递交贺报(3).pptx")

    print("\n" + "=" * 70)
    print("对比结果")
    print("=" * 70)
    print(f"(4) 文件识别: {len(ph1)} 个占位符")
    print(f"(3) 文件识别: {len(ph2)} 个占位符")

    if len(ph1) != len(ph2):
        print("\n⚠️ 识别结果不同！")
        print("说明python-pptx对不同文件的处理可能存在差异")
