#!/usr/bin/env python3
"""
检查WPS设置的Alt Text（替代文本）
WPS可能使用不同的XML属性存储Alt Text
"""

import sys
sys.path.insert(0, 'src')

from pptx import Presentation
from pptx.oxml.ns import qn

pptx_path = r"e:\poster code\poster-again\10月递交贺报(3).pptx"

print("=" * 70)
print("检查WPS设置的Alt Text（替代文本）")
print("=" * 70)
print()

prs = Presentation(pptx_path)

found_alt_text = False

for slide_idx, slide in enumerate(prs.slides, 1):
    print(f"\n--- 幻灯片 {slide_idx} ---")
    slide_has_alt = False

    for shape in slide.shapes:
        shape_name = shape.name

        # 方法1: 标准Alt Text
        alt_text = getattr(shape, 'alt_text', '') or ''

        # 方法2: 检查XML中的descr属性
        descr = ''
        try:
            nvSpPr = shape._element.find(qn('p:nvSpPr'))
            if nvSpPr is not None:
                cNvPr = nvSpPr.find(qn('p:cNvPr'))
                if cNvPr is not None:
                    descr = cNvPr.get('descr', '')
        except:
            pass

        # 方法3: 检查title属性（WPS可能用这个）
        title = ''
        try:
            nvSpPr = shape._element.find(qn('p:nvSpPr'))
            if nvSpPr is not None:
                cNvPr = nvSpPr.find(qn('p:cNvPr'))
                if cNvPr is not None:
                    title = cNvPr.get('title', '')
        except:
            pass

        # 显示所有有内容的
        if alt_text or descr or title:
            found_alt_text = True
            slide_has_alt = True
            print(f"\n  ✓ 形状: {shape_name}")
            if alt_text:
                print(f"    alt_text: '{alt_text}'")
            if descr:
                print(f"    descr: '{descr}'")
            if title:
                print(f"    title: '{title}'")

            # 检查文本内容
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    print(f"    当前文本: '{text[:50]}...'" if len(text) > 50 else f"    当前文本: '{text}'")

    if not slide_has_alt:
        print("  (本页没有设置Alt Text的形状)")

print("\n" + "=" * 70)
if found_alt_text:
    print("✓ 发现设置了Alt Text的形状")
else:
    print("✗ 没有发现任何设置了Alt Text的形状")
print("=" * 70)
