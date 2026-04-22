#!/usr/bin/env python3
"""
检查原始PPTX文件中的Alt Text（替代文本）
"""

import sys
sys.path.insert(0, 'src')

from pptx import Presentation

pptx_path = r"e:\poster code\poster-again\10月递交贺报(3).pptx"

print("=" * 70)
print("检查原始PPTX文件中的Alt Text（替代文本）")
print("=" * 70)
print()

prs = Presentation(pptx_path)

found_alt_text = False

for slide_idx, slide in enumerate(prs.slides, 1):
    print(f"\n--- 幻灯片 {slide_idx} ---")
    slide_has_alt = False

    for shape in slide.shapes:
        shape_name = shape.name

        # 检查 Alt Text
        alt_text = getattr(shape, 'alt_text', '') or ''

        if alt_text:
            found_alt_text = True
            slide_has_alt = True
            print(f"\n  ✓ 形状: {shape_name}")
            print(f"    Alt Text: '{alt_text}'")

            # 检查文本内容
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    print(f"    当前文本: '{text[:50]}...'" if len(text) > 50 else f"    当前文本: '{text}'")

    if not slide_has_alt:
        print("  (本页没有设置Alt Text的形状)")

print("\n" + "=" * 70)
if found_alt_text:
    print("✓ 发现设置了Alt Text的形状")
else:
    print("✗ 没有发现任何设置了Alt Text的形状")
print("=" * 70)
