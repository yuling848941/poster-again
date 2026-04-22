#!/usr/bin/env python3
"""
演示如何正确设置 Alt Text（替代文本）
这个脚本会将形状名称复制到 Alt Text 中
"""

import sys
sys.path.insert(0, 'src')

from pptx import Presentation

pptx_path = r"e:\poster code\poster-again\10月递交贺报(3).pptx"
output_path = r"e:\poster code\poster-again\10月递交贺报_alt_text.pptx"

print("=" * 70)
print("设置 Alt Text（替代文本）演示")
print("=" * 70)
print()

prs = Presentation(pptx_path)

# 定义占位符映射（形状名称 -> Alt Text）
placeholder_mapping = {
    "文本框 7": "ph_标保",
    "文本框 6": "ph_营销员或推广员 姓名",
    "文本框 5": "ph_个险趸期",
    "文本框 4": "ph_支公司 DESC",
    "文本框 3": "ph_险种",
    "文本框 2": "ph_家族 DESC",
}

modified_count = 0

for slide_idx, slide in enumerate(prs.slides, 1):
    print(f"\n--- 幻灯片 {slide_idx} ---")

    for shape in slide.shapes:
        shape_name = shape.name

        # 检查是否需要设置 Alt Text
        if shape_name in placeholder_mapping:
            alt_text_value = placeholder_mapping[shape_name]

            # 设置 Alt Text
            try:
                # 在 python-pptx 中，Alt Text 是通过 shape._element 设置的
                from pptx.oxml.ns import qn

                # 获取或创建 nvSpPr 元素
                nvSpPr = shape._element.find(qn('p:nvSpPr'))
                if nvSpPr is not None:
                    # 查找或创建 cNvPr 元素
                    cNvPr = nvSpPr.find(qn('p:cNvPr'))
                    if cNvPr is not None:
                        # 设置 descr 属性（这就是 Alt Text）
                        cNvPr.set('descr', alt_text_value)
                        modified_count += 1
                        print(f"  ✓ {shape_name} -> Alt Text: '{alt_text_value}'")
                    else:
                        print(f"  ✗ {shape_name} -> 找不到 cNvPr 元素")
                else:
                    print(f"  ✗ {shape_name} -> 找不到 nvSpPr 元素")

            except Exception as e:
                print(f"  ✗ {shape_name} -> 设置失败: {str(e)}")

# 保存文件
prs.save(output_path)

print("\n" + "=" * 70)
print(f"✓ 已修改 {modified_count} 个形状")
print(f"✓ 文件已保存: {output_path}")
print("=" * 70)
print("\n现在你可以用新文件测试 GUI 了！")
