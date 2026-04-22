#!/usr/bin/env python3
"""
调试脚本：查看PPTX文件中的占位符详情
"""

import sys
sys.path.insert(0, 'src')

from pptx import Presentation

pptx_path = r"e:\poster code\poster-again\10月递交贺报(3).pptx"

print("=" * 60)
print("PPTX 文件占位符调试")
print("=" * 60)
print()

prs = Presentation(pptx_path)

for slide_idx, slide in enumerate(prs.slides, 1):
    print(f"\n--- 幻灯片 {slide_idx} ---")

    for shape in slide.shapes:
        print(f"\n  形状: {shape.name}")
        print(f"    类型: {shape.shape_type}")

        # 检查 Alt Text
        alt_text = getattr(shape, 'alt_text', '') or ''
        if alt_text:
            print(f"    Alt Text: '{alt_text}'")

        # 检查是否有文本框
        if shape.has_text_frame:
            text = shape.text_frame.text
            if text:
                print(f"    文本内容: '{text[:100]}...' " if len(text) > 100 else f"    文本内容: '{text}'")

                # 检查文本标记
                if '{{' in text:
                    import re
                    markers = re.findall(r'\{\{(.*?)\}\}', text)
                    if markers:
                        print(f"    发现文本标记: {markers}")

print("\n" + "=" * 60)
print("调试完成")
print("=" * 60)
