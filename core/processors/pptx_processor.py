"""
PPTX 模板处理器 - 纯 Python 实现
基于 python-pptx 库，无需 Microsoft Office 或 WPS

支持占位符识别方式（按优先级）：
1. Alt Text（替代文本）- 推荐，完全隐藏
2. 形状名称（Shape Name）- 备选
3. 文本内容标记（如 {{字段名}}）- 兼容方案
"""

import os
import logging
from typing import Dict, List, Optional, Tuple, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

try:
    from src.core.interfaces.template_processor import (
        ITemplateProcessor,
        ConnectionError,
        TemplateLoadError,
        PlaceholderError,
        SaveError,
    )
except ImportError:
    from core.interfaces.template_processor import (
        ITemplateProcessor,
        ConnectionError,
        TemplateLoadError,
        PlaceholderError,
        SaveError,
    )

logger = logging.getLogger(__name__)


class PPTXProcessor(ITemplateProcessor):
    """
    基于 python-pptx 的模板处理器
    
    无需安装 Microsoft Office 或 WPS，纯 Python 实现
    """
    
    # 占位符前缀
    PLACEHOLDER_PREFIX = "ph_"
    # 文本标记格式
    TEXT_MARKER_START = "{{"
    TEXT_MARKER_END = "}}"
    
    def __init__(self):
        """初始化 PPTX 处理器"""
        super().__init__()
        self.presentation: Optional[Presentation] = None
        self.template_path: Optional[str] = None
        
    def connect_to_application(self) -> bool:
        """
        连接到应用程序
        
        对于 python-pptx，此方法始终返回 True
        因为不需要连接外部应用程序
        
        Returns:
            bool: 始终返回 True
        """
        self.is_connected = True
        logger.info("PPTX 处理器初始化完成（无需外部应用程序）")
        return True
    
    def load_template(self, template_path: str) -> bool:
        """
        加载 PPT 模板文件
        
        Args:
            template_path: 模板文件路径
            
        Returns:
            bool: 加载是否成功
        """
        if not self.validate_template_path(template_path):
            raise TemplateLoadError(f"无效的模板文件路径: {template_path}")
        
        try:
            # 确保路径是绝对路径
            template_path = os.path.abspath(template_path)
            
            # 使用 python-pptx 加载演示文稿
            self.presentation = Presentation(template_path)
            self.template_path = template_path
            
            logger.info(f"成功加载模板: {template_path}")
            logger.info(f"幻灯片数量: {len(self.presentation.slides)}")
            return True
            
        except Exception as e:
            logger.error(f"加载模板失败: {str(e)}")
            raise TemplateLoadError(f"无法加载模板文件: {str(e)}")
    
    def _get_alt_text(self, shape) -> str:
        """
        获取形状的Alt Text（替代文本）
        
        支持多种方式读取，包括WPS设置的Alt Text
        
        Args:
            shape: PPTX形状对象
            
        Returns:
            str: Alt Text内容
        """
        # 方法1: 使用标准alt_text属性
        alt_text = getattr(shape, 'alt_text', '') or ''
        if alt_text:
            return alt_text
        
        # 方法2: 从XML中读取descr属性（WPS使用这种方式）
        try:
            from pptx.oxml.ns import qn
            nvSpPr = shape._element.find(qn('p:nvSpPr'))
            if nvSpPr is not None:
                cNvPr = nvSpPr.find(qn('p:cNvPr'))
                if cNvPr is not None:
                    descr = cNvPr.get('descr', '')
                    if descr:
                        return descr
        except:
            pass
        
        return ''
    
    def find_placeholders(self) -> List[str]:
        """
        查找所有占位符
        
        按以下优先级识别占位符：
        1. Alt Text（替代文本）以 ph_ 开头
        2. 形状名称以 ph_ 开头
        3. 文本内容中的 {{字段名}} 格式
        
        Returns:
            List[str]: 占位符列表（去重）
        """
        if not self.presentation:
            logger.error("未加载演示文稿")
            return []
        
        placeholders = set()
        
        try:
            for slide_idx, slide in enumerate(self.presentation.slides, 1):
                for shape in slide.shapes:
                    # 1. 检查 Alt Text（支持WPS和PowerPoint）
                    alt_text = self._get_alt_text(shape)
                    if alt_text and alt_text.startswith(self.PLACEHOLDER_PREFIX):
                        placeholders.add(alt_text)
                        logger.debug(f"幻灯片 {slide_idx}: 从 Alt Text 发现占位符 '{alt_text}'")
                        continue
                    
                    # 2. 检查形状名称
                    shape_name = getattr(shape, 'name', '') or ''
                    if shape_name and shape_name.startswith(self.PLACEHOLDER_PREFIX):
                        placeholders.add(shape_name)
                        logger.debug(f"幻灯片 {slide_idx}: 从形状名称发现占位符 '{shape_name}'")
                        continue
                    
                    # 3. 检查文本内容中的标记
                    if shape.has_text_frame:
                        text = shape.text_frame.text
                        markers = self._extract_text_markers(text)
                        for marker in markers:
                            placeholders.add(f"{self.PLACEHOLDER_PREFIX}{marker}")
                            logger.debug(f"幻灯片 {slide_idx}: 从文本内容发现占位符 '{marker}'")
            
            result = sorted(list(placeholders))
            logger.info(f"共找到 {len(result)} 个占位符: {result}")
            return result
            
        except Exception as e:
            logger.error(f"查找占位符失败: {str(e)}")
            raise PlaceholderError(f"查找占位符时出错: {str(e)}")
    
    def _extract_text_markers(self, text: str) -> List[str]:
        """
        从文本中提取标记格式的占位符
        
        例如："姓名: {{name}}" -> ["name"]
        
        Args:
            text: 文本内容
            
        Returns:
            List[str]: 提取的占位符名称列表
        """
        markers = []
        start = self.TEXT_MARKER_START
        end = self.TEXT_MARKER_END
        
        idx = 0
        while True:
            start_idx = text.find(start, idx)
            if start_idx == -1:
                break
            end_idx = text.find(end, start_idx + len(start))
            if end_idx == -1:
                break
            
            marker = text[start_idx + len(start):end_idx].strip()
            if marker:
                markers.append(marker)
            idx = end_idx + len(end)
        
        return markers
    
    def get_placeholder_info(self) -> Dict[str, Dict[str, Any]]:
        """
        获取详细的占位符信息
        
        Returns:
            Dict: 占位符详细信息，包括位置、类型等
        """
        if not self.presentation:
            return {}
        
        info = {}
        
        try:
            for slide_idx, slide in enumerate(self.presentation.slides, 1):
                for shape in slide.shapes:
                    placeholder_key = None
                    source = None
                    
                    # 检查 Alt Text（支持WPS和PowerPoint）
                    alt_text = self._get_alt_text(shape)
                    if alt_text and alt_text.startswith(self.PLACEHOLDER_PREFIX):
                        placeholder_key = alt_text
                        source = 'alt_text'
                    
                    # 检查形状名称
                    elif shape.name and shape.name.startswith(self.PLACEHOLDER_PREFIX):
                        placeholder_key = shape.name
                        source = 'shape_name'
                    
                    if placeholder_key:
                        info[placeholder_key] = {
                            'slide': slide_idx,
                            'shape_name': shape.name,
                            'shape_type': str(shape.shape_type),
                            'source': source,
                            'has_text_frame': shape.has_text_frame,
                            'current_text': shape.text_frame.text if shape.has_text_frame else ''
                        }
            
            return info
            
        except Exception as e:
            logger.error(f"获取占位符信息失败: {str(e)}")
            return {}
    
    def replace_placeholders(self, data: Dict[str, str]) -> bool:
        """
        替换占位符内容
        
        支持三种占位符格式：
        1. Alt Text: ph_字段名
        2. 形状名称: ph_字段名
        3. 文本标记: {{字段名}}
        
        Args:
            data: 占位符和对应值的字典
            
        Returns:
            bool: 替换是否成功
        """
        if not self.presentation:
            logger.error("未加载演示文稿")
            return False
        
        try:
            replaced_count = 0
            
            for slide in self.presentation.slides:
                for shape in slide.shapes:
                    # 获取此形状对应的占位符键
                    placeholder_key = self._get_placeholder_key(shape)
                    
                    if placeholder_key and placeholder_key in data:
                        new_value = str(data[placeholder_key])
                        
                        if shape.has_text_frame:
                            self._replace_text_in_shape(shape, new_value)
                            replaced_count += 1
                            logger.debug(f"替换占位符 '{placeholder_key}' -> '{new_value}'")
                    
                    # 处理文本标记格式 {{字段名}}
                    elif shape.has_text_frame:
                        text_replaced = self._replace_text_markers(shape, data)
                        if text_replaced:
                            replaced_count += 1
            
            logger.info(f"成功替换 {replaced_count} 个占位符")
            return True
            
        except Exception as e:
            logger.error(f"替换占位符失败: {str(e)}")
            raise PlaceholderError(f"替换占位符时出错: {str(e)}")
    
    def _get_placeholder_key(self, shape) -> Optional[str]:
        """
        获取形状的占位符键
        
        Args:
            shape: PPT 形状对象
            
        Returns:
            Optional[str]: 占位符键，如果没有则返回 None
        """
        # 优先检查 Alt Text（支持WPS和PowerPoint）
        alt_text = self._get_alt_text(shape)
        if alt_text and alt_text.startswith(self.PLACEHOLDER_PREFIX):
            return alt_text
        
        # 其次检查形状名称
        shape_name = getattr(shape, 'name', '') or ''
        if shape_name and shape_name.startswith(self.PLACEHOLDER_PREFIX):
            return shape_name
        
        return None
    
    def _replace_text_in_shape(self, shape, new_text: str):
        """
        替换形状中的文本，保持格式
        
        Args:
            shape: PPT 形状对象
            new_text: 新文本内容
        """
        if not shape.has_text_frame:
            return
        
        text_frame = shape.text_frame
        
        # 如果只有一个段落，直接替换
        if len(text_frame.paragraphs) == 1:
            paragraph = text_frame.paragraphs[0]
            if len(paragraph.runs) == 1:
                # 只有一个 run，直接替换文本
                paragraph.runs[0].text = new_text
            else:
                # 多个 runs，保留第一个，删除其他的，然后替换文本
                paragraph.runs[0].text = new_text
                for run in paragraph.runs[1:]:
                    run.text = ""
        else:
            # 多个段落，简化处理：清空后设置新文本
            # 这会丢失一些格式，但保持基本样式
            text_frame.text = new_text
    
    def _replace_text_markers(self, shape, data: Dict[str, str]) -> bool:
        """
        替换文本中的标记格式占位符
        
        Args:
            shape: PPT 形状对象
            data: 数据字典
            
        Returns:
            bool: 是否进行了替换
        """
        if not shape.has_text_frame:
            return False
        
        replaced = False
        
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                original_text = run.text
                new_text = original_text
                
                # 替换所有 {{字段名}} 格式的标记
                for key, value in data.items():
                    # 支持 {{字段名}} 和 ph_字段名 两种格式
                    marker1 = f"{self.TEXT_MARKER_START}{key.replace(self.PLACEHOLDER_PREFIX, '')}{self.TEXT_MARKER_END}"
                    marker2 = f"{self.TEXT_MARKER_START}{key}{self.TEXT_MARKER_END}"
                    
                    if marker1 in new_text:
                        new_text = new_text.replace(marker1, str(value))
                        replaced = True
                    if marker2 in new_text:
                        new_text = new_text.replace(marker2, str(value))
                        replaced = True
                
                if new_text != original_text:
                    run.text = new_text
        
        return replaced
    
    def save_presentation(self, output_path: str) -> bool:
        """
        保存演示文稿
        
        Args:
            output_path: 输出文件路径
            
        Returns:
            bool: 保存是否成功
        """
        if not self.presentation:
            logger.error("未加载演示文稿")
            return False
        
        try:
            # 确保输出目录存在
            output_dir = os.path.dirname(output_path)
            if output_dir and not os.path.exists(output_dir):
                os.makedirs(output_dir)
            
            # 保存演示文稿
            self.presentation.save(output_path)
            
            logger.info(f"成功保存演示文稿: {output_path}")
            return True
            
        except Exception as e:
            logger.error(f"保存演示文稿失败: {str(e)}")
            raise SaveError(f"保存演示文稿时出错: {str(e)}")
    
    def close_presentation(self):
        """关闭当前演示文稿"""
        try:
            if self.presentation:
                # python-pptx 不需要显式关闭，只需释放引用
                self.presentation = None
                self.template_path = None
                logger.info("已关闭演示文稿")
        except Exception as e:
            logger.error(f"关闭演示文稿失败: {str(e)}")
    
    def close_application(self):
        """关闭应用程序

        对于 python-pptx，此方法仅释放资源
        """
        self.close_presentation()
        self.is_connected = False
        logger.info("PPTX 处理器已关闭")

    def close(self):
        """关闭处理器，释放资源（别名）"""
        self.close_application()
    
    def get_application_info(self) -> Dict[str, str]:
        """
        获取应用程序信息
        
        Returns:
            Dict: 应用程序信息
        """
        return {
            'name': 'python-pptx',
            'version': '0.6.21+',
            'type': 'pure_python',
            'description': '纯 Python PPT 处理库，无需 Microsoft Office 或 WPS'
        }
    
    def export_to_images(self, output_dir: str, format: str = "JPG", quality: float = 1.0,
                         base_filename: str = None) -> list:
        """
        导出演示文稿为图片

        使用 Microsoft Office 或 WPS Office 导出，需要安装 Office 套件

        Args:
            output_dir: 输出目录
            format: 图片格式 (JPG, PNG)，默认 JPG（Office 兼容性更好）
            quality: 图片质量/分辨率 (1.0-4.0)
            base_filename: 基础文件名（不含扩展名），用于重命名生成的图片

        Returns:
            list: 生成的图片路径列表

        Raises:
            RuntimeError: 如果未安装 Office 套件
        """
        try:
            try:
                from src.core.processors.com_image_exporter import COMImageExporter
            except ImportError:
                from core.processors.com_image_exporter import COMImageExporter

            exporter = COMImageExporter()
            if not exporter.is_available():
                raise RuntimeError(
                    "未检测到 Microsoft Office 或 WPS Office，无法导出图片。\n"
                    "请安装 Office 套件后重试。"
                )

            logger.info(f"使用 {exporter.get_office_info()} 导出图片")

            # 先保存当前演示文稿到临时文件
            import tempfile
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
                temp_path = tmp.name

            self.save_presentation(temp_path)

            try:
                # 转换 quality 参数
                quality_str = "原始大小"
                if quality >= 3.0:
                    quality_str = "3倍"
                elif quality >= 2.0:
                    quality_str = "2倍"

                # 使用 Office 导出
                generated_files = exporter.export_to_images(temp_path, output_dir, format, quality_str)

                # 如果指定了基础文件名，重命名生成的图片
                if base_filename and generated_files:
                    renamed_files = []
                    for i, old_path in enumerate(generated_files, 1):
                        old_name = os.path.basename(old_path)
                        ext = os.path.splitext(old_name)[1]
                        new_name = f"{base_filename}{ext}" if len(generated_files) == 1 else f"{base_filename}_{i}{ext}"
                        new_path = os.path.join(output_dir, new_name)

                        # 如果新文件名已存在，先删除
                        if os.path.exists(new_path) and new_path != old_path:
                            os.remove(new_path)

                        # 重命名文件
                        if old_path != new_path:
                            os.rename(old_path, new_path)
                            logger.info(f"重命名图片: {old_name} -> {new_name}")
                        renamed_files.append(new_path)
                    generated_files = renamed_files

                logger.info(f"成功导出图片到: {output_dir}")
                return generated_files
            finally:
                # 清理临时文件
                try:
                    os.remove(temp_path)
                except:
                    pass

        except RuntimeError:
            raise
        except Exception as e:
            logger.error(f"导出图片时出错: {e}")
            return []

    def get_slide_count(self) -> int:
        """
        获取幻灯片数量

        Returns:
            int: 幻灯片数量
        """
        if self.presentation:
            return len(self.presentation.slides)
        return 0

    def get_shape_count(self) -> int:
        """
        获取形状总数

        Returns:
            int: 形状总数
        """
        if not self.presentation:
            return 0

        count = 0
        for slide in self.presentation.slides:
            count += len(slide.shapes)
        return count
